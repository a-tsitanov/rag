"""Document parser supporting PDF, DOCX, PPTX, TXT/MD, and EML formats.

Each format extractor returns as much data as it can.  On partial failure
the parser logs the error and returns whatever it managed to extract —
callers always get a ``ParsedDocument``, never an exception.
"""

from __future__ import annotations

import email
import logging
import mimetypes
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path

import chardet

logger = logging.getLogger(__name__)

# ── data classes ──────────────────────────────────────────────────────


@dataclass
class Section:
    title: str
    content: str
    level: int = 0  # heading depth (1 = H1, 0 = body)


@dataclass
class ParsedDocument:
    text: str
    metadata: dict = field(default_factory=dict)
    sections: list[Section] = field(default_factory=list)


# ── format detection ──────────────────────────────────────────────────

_EXT_TO_TYPE: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt": "txt",
    ".md": "md",
    ".markdown": "md",
    ".eml": "eml",
}


def _detect_type(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in _EXT_TO_TYPE:
        return _EXT_TO_TYPE[ext]
    mime, _ = mimetypes.guess_type(str(path))
    if mime:
        if "pdf" in mime:
            return "pdf"
        if "wordprocessing" in mime or "msword" in mime:
            return "docx"
        if "presentation" in mime:
            return "pptx"
        if "message/rfc822" in mime:
            return "eml"
    return "txt"


# ── format parsers ────────────────────────────────────────────────────


def _parse_pdf(path: Path) -> ParsedDocument:
    pages_text: list[str] = []
    metadata: dict = {}

    # primary: pdfplumber (better table / layout handling)
    try:
        import pdfplumber

        with pdfplumber.open(path) as pdf:
            metadata["pages"] = len(pdf.pages)
            if pdf.metadata:
                metadata["title"] = pdf.metadata.get("Title", "")
                metadata["author"] = pdf.metadata.get("Author", "")
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text)
    except Exception as exc:
        logger.warning("pdfplumber failed for %s, falling back to pypdf: %s", path, exc)
        pages_text.clear()

    # fallback: pypdf
    if not pages_text:
        try:
            from pypdf import PdfReader

            reader = PdfReader(path)
            metadata["pages"] = len(reader.pages)
            info = reader.metadata
            if info:
                metadata["title"] = info.get("/Title", "")
                metadata["author"] = info.get("/Author", "")
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text)
        except Exception as exc:
            logger.error("pypdf also failed for %s: %s", path, exc)
            metadata["error"] = str(exc)

    full_text = "\n\n".join(pages_text)
    sections = [
        Section(title=f"Page {i + 1}", content=t, level=0)
        for i, t in enumerate(pages_text)
    ]

    return ParsedDocument(text=full_text, metadata=metadata, sections=sections)


def _parse_docx(path: Path) -> ParsedDocument:
    from docx import Document as DocxDocument

    doc = DocxDocument(path)
    metadata: dict = {}
    core = doc.core_properties
    if core:
        metadata["title"] = core.title or ""
        metadata["author"] = core.author or ""
        if core.created:
            metadata["created_at"] = core.created.isoformat()

    sections: list[Section] = []
    current_body: list[str] = []
    current_heading = ""
    current_level = 0

    for para in doc.paragraphs:
        style_name = (para.style.name or "").lower()
        if style_name.startswith("heading"):
            # flush previous section
            if current_body or current_heading:
                sections.append(Section(
                    title=current_heading,
                    content="\n".join(current_body),
                    level=current_level,
                ))
                current_body = []
            current_heading = para.text
            # extract level: "Heading 2" -> 2
            try:
                current_level = int(style_name.split()[-1])
            except (ValueError, IndexError):
                current_level = 1
        else:
            if para.text.strip():
                current_body.append(para.text)

    # flush last section
    if current_body or current_heading:
        sections.append(Section(
            title=current_heading,
            content="\n".join(current_body),
            level=current_level,
        ))

    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return ParsedDocument(text=full_text, metadata=metadata, sections=sections)


def _parse_pptx(path: Path) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(path)
    metadata: dict = {}
    core = prs.core_properties
    if core:
        metadata["title"] = core.title or ""
        metadata["author"] = core.author or ""

    sections: list[Section] = []
    all_text: list[str] = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        slide_title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_parts.append(text)
                    if not slide_title and shape.shape_id == slide.shapes[0].shape_id:
                        slide_title = text

        content = "\n".join(slide_parts)
        if content:
            sections.append(Section(
                title=slide_title or f"Slide {idx}",
                content=content,
                level=0,
            ))
            all_text.append(content)

    metadata["pages"] = len(prs.slides)
    full_text = "\n\n".join(all_text)
    return ParsedDocument(text=full_text, metadata=metadata, sections=sections)


def _parse_text(path: Path) -> ParsedDocument:
    raw = path.read_bytes()
    detected = chardet.detect(raw)
    encoding = detected.get("encoding") or "utf-8"

    try:
        text = raw.decode(encoding)
    except (UnicodeDecodeError, LookupError):
        text = raw.decode("utf-8", errors="replace")

    sections = [Section(title=path.name, content=text, level=0)]
    metadata: dict = {"encoding": encoding}
    return ParsedDocument(text=text, metadata=metadata, sections=sections)


def _parse_eml(path: Path) -> ParsedDocument:
    raw = path.read_bytes()
    msg = email.message_from_bytes(raw)

    metadata: dict = {
        "title": msg.get("Subject", ""),
        "author": msg.get("From", ""),
    }
    date_str = msg.get("Date", "")
    if date_str:
        metadata["created_at"] = date_str

    # body
    body_parts: list[str] = []
    attachments: list[str] = []

    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            disp = str(part.get("Content-Disposition", ""))
            if "attachment" in disp:
                fname = part.get_filename() or "unnamed"
                attachments.append(fname)
                continue
            if ct == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    body_parts.append(payload.decode("utf-8", errors="replace"))
            elif ct == "text/html" and not body_parts:
                from bs4 import BeautifulSoup

                payload = part.get_payload(decode=True)
                if payload:
                    soup = BeautifulSoup(
                        payload.decode("utf-8", errors="replace"), "html.parser"
                    )
                    body_parts.append(soup.get_text(separator="\n"))
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            body_parts.append(payload.decode("utf-8", errors="replace"))

    if attachments:
        metadata["attachments"] = attachments

    text = "\n\n".join(body_parts)
    sections = []
    subject = metadata.get("title", "")
    if subject:
        sections.append(Section(title="Subject", content=subject, level=1))
    if text:
        sections.append(Section(title="Body", content=text, level=0))
    if attachments:
        sections.append(Section(
            title="Attachments",
            content="\n".join(attachments),
            level=0,
        ))

    return ParsedDocument(text=text, metadata=metadata, sections=sections)


# ── dispatcher ────────────────────────────────────────────────────────

_PARSERS: dict[str, callable] = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "txt": _parse_text,
    "md": _parse_text,
    "eml": _parse_eml,
}


def _enrich_metadata(doc: ParsedDocument, path: Path, doc_type: str):
    doc.metadata.setdefault("doc_type", doc_type)
    doc.metadata.setdefault("title", path.stem)
    doc.metadata["word_count"] = len(doc.text.split()) if doc.text else 0

    try:
        stat = path.stat()
        if "created_at" not in doc.metadata:
            ts = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)
            doc.metadata["created_at"] = ts.isoformat()
    except OSError:
        pass


# ── public API ────────────────────────────────────────────────────────


class DocumentParser:
    """Parse documents of supported formats into a uniform structure."""

    def parse(self, path: Path) -> ParsedDocument:
        path = Path(path)
        doc_type = _detect_type(path)
        parser_fn = _PARSERS.get(doc_type)

        if parser_fn is None:
            logger.warning("Unsupported format %r, treating as plain text", doc_type)
            parser_fn = _parse_text

        try:
            doc = parser_fn(path)
        except Exception as exc:
            logger.error("Failed to parse %s: %s", path, exc, exc_info=True)
            doc = ParsedDocument(text="", metadata={"error": str(exc)})

        _enrich_metadata(doc, path, doc_type)
        return doc
